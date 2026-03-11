"""
PPTX report renderer using python-pptx.

Generates professional, enterprise-grade PowerPoint presentations with charts,
tables, and branded styling for the KhushFus social listening platform.
"""

import logging
from datetime import datetime
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTXRenderer:
    """Renders report data into a branded PPTX presentation."""

    # Brand colors
    PRIMARY = (37, 99, 235)
    DARK_BG = (15, 23, 42)
    CARD_BG = (30, 41, 59)
    WHITE = (248, 250, 252)
    MUTED = (148, 163, 184)
    GREEN = (16, 185, 129)
    RED = (239, 68, 68)
    YELLOW = (245, 158, 11)

    def __init__(self, report_data: dict, report_type: str, project_name: str):
        self.data = report_data
        self.report_type = report_type
        self.project_name = project_name

    def render(self, output_path: str | None = None) -> str:
        """Generate the PPTX file and return the file path."""
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        def rgb(color_tuple):
            return RGBColor(*color_tuple)

        def add_bg(slide, color=self.DARK_BG):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = rgb(color)

        def add_text(slide, left, top, width, height, text, size=14, color=self.WHITE, bold=False, align=PP_ALIGN.LEFT):
            tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(text)
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(color)
            p.font.bold = bold
            p.alignment = align
            return tx_box

        def add_bar(slide, x, y, width, height, color):
            bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(max(width, 0.1)), Inches(height))
            bar.fill.solid()
            bar.fill.fore_color.rgb = rgb(color)
            bar.line.fill.background()
            return bar

        def add_card(slide, x, y, w, h, label, value):
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(self.CARD_BG)
            shape.line.fill.background()
            add_text(slide, x + 0.3, y + 0.3, w - 0.6, 0.5, label, 14, self.MUTED)
            add_text(slide, x + 0.3, y + 0.9, w - 0.6, 0.8, value, 32, self.WHITE, bold=True)

        # --- Slide 1: Title ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text(slide, 1, 1.5, 11, 1, "KhushFus", 18, self.MUTED, bold=True)
        add_text(slide, 1, 2.5, 11, 1.5, f"{self.report_type.capitalize()} Report", 44, self.WHITE, bold=True)
        add_text(slide, 1, 4.2, 11, 0.8, self.project_name, 28, self.PRIMARY, bold=True)

        period = self.data.get("period", {})
        period_str = f"{period.get('start', '')[:10]}  —  {period.get('end', '')[:10]}"
        add_text(slide, 1, 5.3, 11, 0.5, period_str, 16, self.MUTED)
        gen_time = datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')
        add_text(slide, 1, 6.2, 11, 0.5, f"Generated: {gen_time}", 12, self.MUTED)

        # --- Slide 2: Executive Summary ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text(slide, 0.8, 0.4, 6, 0.7, "Executive Summary", 28, self.WHITE, bold=True)

        engagement = self.data.get("engagement", {})
        kpis = [
            ("Total Mentions", f"{self.data.get('total_mentions', 0):,}"),
            ("Total Reach", f"{engagement.get('total_reach', 0):,}"),
            ("Total Likes", f"{engagement.get('total_likes', 0):,}"),
            ("Total Shares", f"{engagement.get('total_shares', 0):,}"),
            ("Total Comments", f"{engagement.get('total_comments', 0):,}"),
            ("Flagged", f"{self.data.get('flagged_mentions', 0):,}"),
        ]
        for i, (label, value) in enumerate(kpis):
            col, row = i % 3, i // 3
            add_card(slide, 0.8 + col * 4, 1.5 + row * 2.5, 3.6, 2, label, value)

        # --- Slide 3: Sentiment ---
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text(slide, 0.8, 0.4, 6, 0.7, "Sentiment Analysis", 28, self.WHITE, bold=True)

        sentiment = self.data.get("sentiment", {})
        breakdown = sentiment.get("breakdown", {})
        avg_score = sentiment.get("average_score", 0)
        total_s = max(sum(breakdown.values()), 1)

        color_map = {"positive": self.GREEN, "negative": self.RED, "neutral": self.MUTED, "mixed": self.YELLOW}
        y = 1.5
        for label, count in breakdown.items():
            pct = round(count / total_s * 100, 1)
            c = color_map.get(label, self.MUTED)
            add_text(slide, 0.8, y, 2.5, 0.5, f"{label.capitalize()}: {count} ({pct}%)", 16, c, bold=True)
            add_bar(slide, 3.5, y + 0.05, 8, 0.4, (51, 65, 85))  # bg
            add_bar(slide, 3.5, y + 0.05, pct / 100 * 8, 0.4, c)
            y += 0.8

        add_text(slide, 0.8, y + 0.5, 5, 0.5, f"Average Sentiment Score: {avg_score}", 18, self.WHITE, bold=True)

        # --- Slide 4: Platforms ---
        platforms = self.data.get("platforms", {})
        if platforms:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text(slide, 0.8, 0.4, 6, 0.7, "Platform Breakdown", 28, self.WHITE, bold=True)

            ptotal = max(sum(platforms.values()), 1)
            y = 1.5
            for plat, count in sorted(platforms.items(), key=lambda x: x[1], reverse=True):
                pct = round(count / ptotal * 100, 1)
                add_text(slide, 0.8, y, 3, 0.5, f"{plat}: {count} ({pct}%)", 15, self.WHITE)
                add_bar(slide, 4.5, y + 0.05, pct / 100 * 7, 0.35, self.PRIMARY)
                y += 0.7

        # --- Slide 5: Top Contributors ---
        contributors = self.data.get("top_contributors", [])
        if contributors:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text(slide, 0.8, 0.4, 6, 0.7, "Top Contributors", 28, self.WHITE, bold=True)

            headers = ["Name", "Handle", "Platform", "Followers", "Mentions"]
            col_widths = [3, 3, 2, 2, 1.5]
            y = 1.5
            x = 0.8
            for h, w in zip(headers, col_widths):
                add_text(slide, x, y, w, 0.4, h, 12, self.PRIMARY, bold=True)
                x += w
            y += 0.5

            for c in contributors[:12]:
                x = 0.8
                vals = [str(c.get("name", "")), str(c.get("handle", "")), str(c.get("platform", "")),
                        f"{c.get('followers', 0):,}", str(c.get("mentions", 0))]
                for v, w in zip(vals, col_widths):
                    add_text(slide, x, y, w, 0.4, v, 11, self.WHITE)
                    x += w
                y += 0.45

        # --- Slide 6: Keywords ---
        keywords = self.data.get("keyword_frequency", {})
        if keywords:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text(slide, 0.8, 0.4, 6, 0.7, "Keyword Performance", 28, self.WHITE, bold=True)

            sorted_kw = sorted(keywords.items(), key=lambda x: x[1], reverse=True)[:15]
            kw_max = max(v for _, v in sorted_kw) if sorted_kw else 1
            y = 1.5
            for kw, freq in sorted_kw:
                add_text(slide, 0.8, y, 3, 0.4, f"{kw}: {freq}", 13, self.WHITE)
                add_bar(slide, 4.5, y + 0.05, freq / kw_max * 7, 0.3, self.GREEN)
                y += 0.55

        # --- Slide 7: Influencers ---
        influencers = self.data.get("influencers", [])
        if influencers:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text(slide, 0.8, 0.4, 6, 0.7, "Key Influencers", 28, self.WHITE, bold=True)

            headers = ["Name", "Handle", "Platform", "Followers"]
            col_widths = [3.5, 3.5, 2.5, 2]
            y = 1.5
            x = 0.8
            for h, w in zip(headers, col_widths):
                add_text(slide, x, y, w, 0.4, h, 12, self.PRIMARY, bold=True)
                x += w
            y += 0.5

            for inf in influencers[:10]:
                x = 0.8
                vals = [str(inf.get("name", "")), str(inf.get("handle", "")),
                        str(inf.get("platform", "")), f"{inf.get('followers', 0):,}"]
                for v, w in zip(vals, col_widths):
                    add_text(slide, x, y, w, 0.4, v, 11, self.WHITE)
                    x += w
                y += 0.45

        # Save
        if not output_path:
            output_dir = Path("reports_output")
            output_dir.mkdir(exist_ok=True)
            ts = datetime.utcnow().strftime('%Y%m%d_%H%M%S')
            output_path = str(output_dir / f"{self.project_name}_{self.report_type}_{ts}.pptx")

        prs.save(output_path)
        logger.info(f"PPTX report saved: {output_path}")
        return output_path
