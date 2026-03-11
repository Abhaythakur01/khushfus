"""
Report Service — generates hourly/daily/weekly/monthly/quarterly/yearly reports.

1. Listens to 'reports:request' stream for on-demand generation
2. Runs a scheduler for automatic report generation
3. Aggregates mention data from PostgreSQL
4. Generates PDF (WeasyPrint + Jinja2) or PPTX (python-pptx) reports
5. Publishes 'reports:ready' event when done
"""

import asyncio
import json
import logging
import os
import signal
from datetime import datetime, timedelta
from pathlib import Path

from jinja2 import Environment, FileSystemLoader, TemplateNotFound, TemplateSyntaxError
from sqlalchemy import func, select

from shared.database import create_db
from shared.events import STREAM_REPORT_READY, STREAM_REPORT_REQUESTS, EventBus
from shared.logging_config import setup_logging
from shared.models import Mention, Platform, Project, ProjectStatus, Report, ReportType, Sentiment
from shared.service_utils import ConsumerMetrics, backoff_with_jitter, is_transient_error, start_liveness_server

setup_logging("report-service")
logger = logging.getLogger(__name__)

DATABASE_URL = os.getenv("DATABASE_URL", "postgresql+asyncpg://khushfus:khushfus_dev@postgres:5432/khushfus")
REDIS_URL = os.getenv("REDIS_URL", "redis://redis:6379/0")

GROUP_NAME = "report-service"
CONSUMER_NAME = f"report-{os.getpid()}"
TEMPLATE_DIR = Path(__file__).parent / "templates"
OUTPUT_DIR = Path("/app/reports_output")

metrics = ConsumerMetrics("report-service")

PERIOD_MAP = {
    "hourly": timedelta(hours=1),
    "daily": timedelta(days=1),
    "weekly": timedelta(weeks=1),
    "monthly": timedelta(days=30),
    "quarterly": timedelta(days=90),
    "yearly": timedelta(days=365),
}


async def build_report_data(db, project_id: int, start: datetime, end: datetime) -> dict:
    """Aggregate mention data for a report period."""
    base = [Mention.project_id == project_id, Mention.collected_at >= start, Mention.collected_at <= end]

    total = (await db.execute(select(func.count(Mention.id)).where(*base))).scalar() or 0

    sentiment_counts = {}
    for s in Sentiment:
        c = (await db.execute(select(func.count(Mention.id)).where(*base, Mention.sentiment == s))).scalar() or 0
        sentiment_counts[s.value] = c

    avg = (await db.execute(select(func.avg(Mention.sentiment_score)).where(*base))).scalar() or 0.0

    platform_counts = {}
    for p in Platform:
        c = (await db.execute(select(func.count(Mention.id)).where(*base, Mention.platform == p))).scalar() or 0
        if c > 0:
            platform_counts[p.value] = c

    eng = (
        await db.execute(
            select(
                func.sum(Mention.likes), func.sum(Mention.shares), func.sum(Mention.comments), func.sum(Mention.reach)
            ).where(*base)
        )
    ).one()

    top = await db.execute(
        select(
            Mention.author_name,
            Mention.author_handle,
            Mention.author_followers,
            Mention.platform,
            func.count(Mention.id).label("cnt"),
        )
        .where(*base)
        .group_by(Mention.author_handle, Mention.author_name, Mention.author_followers, Mention.platform)
        .order_by(func.count(Mention.id).desc())
        .limit(20)
    )
    top_contributors = [
        {
            "name": r.author_name,
            "handle": r.author_handle,
            "followers": r.author_followers,
            "platform": r.platform,
            "mentions": r.cnt,
        }
        for r in top
    ]

    # Keyword frequency
    kw_result = await db.execute(select(Mention.matched_keywords).where(*base, Mention.matched_keywords.isnot(None)))
    keyword_freq: dict[str, int] = {}
    for row in kw_result.scalars():
        for kw in row.split(","):
            kw = kw.strip()
            if kw:
                keyword_freq[kw] = keyword_freq.get(kw, 0) + 1

    flagged = (
        await db.execute(select(func.count(Mention.id)).where(*base, Mention.is_flagged.is_(True)))
    ).scalar() or 0

    influencers = await db.execute(
        select(
            Mention.author_name,
            Mention.author_handle,
            Mention.author_followers,
            Mention.author_profile_url,
            Mention.platform,
        )
        .where(*base, Mention.author_followers > 1000)
        .distinct(Mention.author_handle)
        .order_by(Mention.author_handle, Mention.author_followers.desc())
        .limit(10)
    )

    return {
        "period": {"start": start.isoformat(), "end": end.isoformat()},
        "total_mentions": total,
        "sentiment": {"breakdown": sentiment_counts, "average_score": round(float(avg), 4)},
        "platforms": platform_counts,
        "top_contributors": top_contributors,
        "engagement": {
            "total_likes": eng[0] or 0,
            "total_shares": eng[1] or 0,
            "total_comments": eng[2] or 0,
            "total_reach": eng[3] or 0,
        },
        "keyword_frequency": keyword_freq,
        "flagged_mentions": flagged,
        "influencers": [
            {
                "name": r.author_name,
                "handle": r.author_handle,
                "followers": r.author_followers,
                "profile_url": r.author_profile_url,
                "platform": r.platform,
            }
            for r in influencers
        ],
    }


async def generate_report(session_factory, project_id: int, report_type: str, fmt: str = "pdf") -> int:
    """Generate a report, store it, create PDF or PPTX. Returns report ID."""
    now = datetime.utcnow()
    period_start = now - PERIOD_MAP.get(report_type, timedelta(days=1))

    async with session_factory() as db:
        project = await db.get(Project, project_id)
        if not project:
            logger.error(f"Project {project_id} not found")
            return 0

        data = await build_report_data(db, project_id, period_start, now)

        report = Report(
            project_id=project_id,
            report_type=ReportType(report_type),
            title=f"{report_type.capitalize()} Report - {now.strftime('%Y-%m-%d %H:%M')}",
            format=fmt,
            period_start=period_start,
            period_end=now,
            data_json=json.dumps(data, default=str),
        )

        if fmt == "pptx":
            file_path = generate_pptx(data, report_type, project.name)
        else:
            file_path = generate_pdf(data, report_type, project.name)

        if file_path:
            report.file_path = file_path

        db.add(report)
        await db.commit()
        await db.refresh(report)

        logger.info(f"Generated {report_type} {fmt.upper()} report #{report.id} for project {project_id}")
        return report.id


def generate_pdf(report_data: dict, report_type: str, project_name: str) -> str | None:
    """Render HTML template and convert to PDF."""
    OUTPUT_DIR.mkdir(exist_ok=True)
    filepath = None

    try:
        env = Environment(loader=FileSystemLoader(str(TEMPLATE_DIR)))
        template = env.get_template("report.html")
        html = template.render(
            report=report_data,
            report_type=report_type,
            project_name=project_name,
            generated_at=datetime.utcnow().strftime("%Y-%m-%d %H:%M UTC"),
        )

        filename = f"{project_name}_{report_type}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = OUTPUT_DIR / filename

        try:
            from weasyprint import HTML

            HTML(string=html).write_pdf(str(filepath))
            return str(filepath)
        except ImportError:
            html_path = filepath.with_suffix(".html")
            html_path.write_text(html, encoding="utf-8")
            return str(html_path)
    except TemplateNotFound as e:
        logger.error("Report template not found: %s — check TEMPLATE_DIR=%s", e, TEMPLATE_DIR)
        return None
    except TemplateSyntaxError as e:
        logger.error("Report template syntax error in %s line %s: %s", e.filename, e.lineno, e.message)
        return None
    except Exception as e:
        if filepath and filepath.exists():
            filepath.unlink(missing_ok=True)
        logger.error(f"PDF generation failed: {e}")
        return None


def generate_pptx(report_data: dict, report_type: str, project_name: str) -> str | None:
    """Generate a PPTX presentation report."""
    OUTPUT_DIR.mkdir(exist_ok=True)
    filepath = None

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # -- Color palette --
        primary = RGBColor(37, 99, 235)
        dark_bg = RGBColor(15, 23, 42)
        card_bg = RGBColor(30, 41, 59)
        white = RGBColor(248, 250, 252)
        muted = RGBColor(148, 163, 184)
        green = RGBColor(16, 185, 129)
        red = RGBColor(239, 68, 68)
        yellow = RGBColor(245, 158, 11)

        def add_bg(slide, color=dark_bg):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def add_text_box(
            slide, left, top, width, height, text,
            font_size=14, color=white, bold=False, align=PP_ALIGN.LEFT,
        ):
            tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = tx_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = str(text)
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.alignment = align
            return tx_box

        # ===== SLIDE 1: Title =====
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        add_bg(slide)
        add_text_box(slide, 1, 1.5, 11, 1, "KhushFus", 18, muted, bold=True)
        add_text_box(
            slide, 1, 2.5, 11, 1.5, f"{report_type.capitalize()} Report",
            44, white, bold=True, align=PP_ALIGN.LEFT,
        )
        add_text_box(slide, 1, 4.2, 11, 0.8, project_name, 28, primary, bold=True)
        period = report_data.get("period", {})
        period_str = f"{period.get('start', '')[:10]}  —  {period.get('end', '')[:10]}"
        add_text_box(slide, 1, 5.3, 11, 0.5, period_str, 16, muted)
        gen_str = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"
        add_text_box(slide, 1, 6.2, 11, 0.5, gen_str, 12, muted)

        # ===== SLIDE 2: Executive Summary =====
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text_box(slide, 0.8, 0.4, 6, 0.7, "Executive Summary", 28, white, bold=True)

        engagement = report_data.get("engagement", {})
        kpis = [
            ("Total Mentions", f"{report_data.get('total_mentions', 0):,}"),
            ("Total Reach", f"{engagement.get('total_reach', 0):,}"),
            ("Total Likes", f"{engagement.get('total_likes', 0):,}"),
            ("Total Shares", f"{engagement.get('total_shares', 0):,}"),
            ("Total Comments", f"{engagement.get('total_comments', 0):,}"),
            ("Flagged", f"{report_data.get('flagged_mentions', 0):,}"),
        ]
        for i, (label, value) in enumerate(kpis):
            col = i % 3
            row = i // 3
            x = 0.8 + col * 4
            y = 1.5 + row * 2.5

            # Card background
            shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.6), Inches(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = card_bg
            shape.line.fill.background()

            add_text_box(slide, x + 0.3, y + 0.3, 3, 0.5, label, 14, muted)
            add_text_box(slide, x + 0.3, y + 0.9, 3, 0.8, value, 32, white, bold=True)

        # ===== SLIDE 3: Sentiment Analysis =====
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        add_text_box(slide, 0.8, 0.4, 6, 0.7, "Sentiment Analysis", 28, white, bold=True)

        sentiment = report_data.get("sentiment", {})
        breakdown = sentiment.get("breakdown", {})
        avg_score = sentiment.get("average_score", 0)
        total = max(sum(breakdown.values()), 1)

        sentiment_colors = {"positive": green, "negative": red, "neutral": muted, "mixed": yellow}
        y_pos = 1.5
        for label, count in breakdown.items():
            pct = round(count / total * 100, 1)
            color = sentiment_colors.get(label, muted)
            add_text_box(slide, 0.8, y_pos, 2.5, 0.5, f"{label.capitalize()}: {count} ({pct}%)", 16, color, bold=True)
            # Bar background
            bar_bg = slide.shapes.add_shape(1, Inches(3.5), Inches(y_pos + 0.05), Inches(8), Inches(0.4))
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = RGBColor(51, 65, 85)
            bar_bg.line.fill.background()
            # Bar fill
            bar_w = max(pct / 100 * 8, 0.1)
            bar = slide.shapes.add_shape(1, Inches(3.5), Inches(y_pos + 0.05), Inches(bar_w), Inches(0.4))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            y_pos += 0.8

        add_text_box(slide, 0.8, y_pos + 0.5, 5, 0.5, f"Average Sentiment Score: {avg_score}", 18, white, bold=True)

        # ===== SLIDE 4: Platform Breakdown =====
        platforms = report_data.get("platforms", {})
        if platforms:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text_box(slide, 0.8, 0.4, 6, 0.7, "Platform Breakdown", 28, white, bold=True)

            platform_total = max(sum(platforms.values()), 1)
            y_pos = 1.5
            for plat, count in sorted(platforms.items(), key=lambda x: x[1], reverse=True):
                pct = round(count / platform_total * 100, 1)
                add_text_box(slide, 0.8, y_pos, 3, 0.5, f"{plat}: {count} ({pct}%)", 15, white)
                bar_w = max(pct / 100 * 7, 0.1)
                bar = slide.shapes.add_shape(1, Inches(4.5), Inches(y_pos + 0.05), Inches(bar_w), Inches(0.35))
                bar.fill.solid()
                bar.fill.fore_color.rgb = primary
                bar.line.fill.background()
                y_pos += 0.7

        # ===== SLIDE 5: Top Contributors =====
        contributors = report_data.get("top_contributors", [])
        if contributors:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text_box(slide, 0.8, 0.4, 6, 0.7, "Top Contributors", 28, white, bold=True)

            headers = ["Name", "Handle", "Platform", "Followers", "Mentions"]
            col_widths = [3, 3, 2, 2, 1.5]
            x_start = 0.8
            y_pos = 1.5

            # Header row
            x = x_start
            for h, w in zip(headers, col_widths):
                add_text_box(slide, x, y_pos, w, 0.4, h, 12, primary, bold=True)
                x += w
            y_pos += 0.5

            for c in contributors[:12]:
                x = x_start
                vals = [
                    str(c.get("name", "")),
                    str(c.get("handle", "")),
                    str(c.get("platform", "")),
                    f"{c.get('followers', 0):,}",
                    str(c.get("mentions", 0)),
                ]
                for v, w in zip(vals, col_widths):
                    add_text_box(slide, x, y_pos, w, 0.4, v, 11, white)
                    x += w
                y_pos += 0.45

        # ===== SLIDE 6: Keyword Performance =====
        keywords = report_data.get("keyword_frequency", {})
        if keywords:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text_box(slide, 0.8, 0.4, 6, 0.7, "Keyword Performance", 28, white, bold=True)

            sorted_kw = sorted(keywords.items(), key=lambda x: x[1], reverse=True)[:15]
            kw_max = max(v for _, v in sorted_kw) if sorted_kw else 1
            y_pos = 1.5
            for kw, freq in sorted_kw:
                add_text_box(slide, 0.8, y_pos, 3, 0.4, f"{kw}: {freq}", 13, white)
                bar_w = max(freq / kw_max * 7, 0.1)
                bar = slide.shapes.add_shape(1, Inches(4.5), Inches(y_pos + 0.05), Inches(bar_w), Inches(0.3))
                bar.fill.solid()
                bar.fill.fore_color.rgb = green
                bar.line.fill.background()
                y_pos += 0.55

        # ===== SLIDE 7: Key Influencers =====
        influencers = report_data.get("influencers", [])
        if influencers:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide)
            add_text_box(slide, 0.8, 0.4, 6, 0.7, "Key Influencers", 28, white, bold=True)

            headers = ["Name", "Handle", "Platform", "Followers"]
            col_widths = [3.5, 3.5, 2.5, 2]
            x_start = 0.8
            y_pos = 1.5

            x = x_start
            for h, w in zip(headers, col_widths):
                add_text_box(slide, x, y_pos, w, 0.4, h, 12, primary, bold=True)
                x += w
            y_pos += 0.5

            for inf in influencers[:10]:
                x = x_start
                vals = [
                    str(inf.get("name", "")),
                    str(inf.get("handle", "")),
                    str(inf.get("platform", "")),
                    f"{inf.get('followers', 0):,}",
                ]
                for v, w in zip(vals, col_widths):
                    add_text_box(slide, x, y_pos, w, 0.4, v, 11, white)
                    x += w
                y_pos += 0.45

        # Save
        filename = f"{project_name}_{report_type}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = OUTPUT_DIR / filename
        prs.save(str(filepath))
        logger.info(f"PPTX report saved: {filepath}")
        return str(filepath)

    except ImportError:
        logger.warning("python-pptx not installed, falling back to PDF")
        return generate_pdf(report_data, report_type, project_name)
    except Exception as e:
        if filepath and filepath.exists():
            filepath.unlink(missing_ok=True)
        logger.error(f"PPTX generation failed: {e}")
        return None


async def listen_for_requests(bus: EventBus, session_factory, shutdown_event: asyncio.Event):
    """Listen for on-demand report generation requests."""
    await bus.ensure_group(STREAM_REPORT_REQUESTS, GROUP_NAME)
    logger.info("Listening for report requests...")
    retry_attempt = 0

    while not shutdown_event.is_set():
        try:
            messages = await bus.consume(STREAM_REPORT_REQUESTS, GROUP_NAME, CONSUMER_NAME, block_ms=3000)
            retry_attempt = 0
            for msg_id, data in messages:
                project_id = int(data.get("project_id", 0))
                report_type = data.get("report_type", "daily")
                fmt = data.get("format", "pdf")
                try:
                    report_id = await generate_report(session_factory, project_id, report_type, fmt)
                    if report_id:
                        await bus.publish(
                            STREAM_REPORT_READY,
                            {
                                "project_id": project_id,
                                "report_id": report_id,
                                "report_type": report_type,
                                "format": fmt,
                            },
                        )
                    metrics.record_processed()
                except Exception as e:
                    metrics.record_failed()
                    if is_transient_error(e):
                        delay = backoff_with_jitter(retry_attempt)
                        logger.error(
                            "Transient error processing report %s, retry in %.1fs: %s", msg_id, delay, e
                        )
                        await asyncio.sleep(delay)
                        retry_attempt += 1
                        continue
                    else:
                        logger.error(f"Permanent error processing report request {msg_id}, acking: {e}")
                await bus.ack(STREAM_REPORT_REQUESTS, GROUP_NAME, msg_id)
        except asyncio.CancelledError:
            break
        except Exception as e:
            logger.error(f"Report request processing error: {e}")
            await asyncio.sleep(1)


async def scheduled_reports(bus: EventBus, session_factory, shutdown_event: asyncio.Event):
    """Generate reports on schedule.

    - Hourly: every hour at :00
    - Daily: 6:00 AM UTC
    - Weekly: Monday 7:00 AM UTC
    - Monthly: 1st of month 8:00 AM UTC
    - Quarterly: 1st of Jan/Apr/Jul/Oct 9:00 AM UTC
    - Yearly: January 1st 10:00 AM UTC
    """
    while not shutdown_event.is_set():
        now = datetime.utcnow()

        tasks = []
        if now.minute < 5:
            tasks.append("hourly")
        if now.hour == 6 and now.minute < 5:
            tasks.append("daily")
        if now.hour == 7 and now.minute < 5 and now.weekday() == 0:
            tasks.append("weekly")
        if now.hour == 8 and now.minute < 5 and now.day == 1:
            tasks.append("monthly")
        if now.hour == 9 and now.minute < 5 and now.day == 1 and now.month in (1, 4, 7, 10):
            tasks.append("quarterly")
        if now.hour == 10 and now.minute < 5 and now.day == 1 and now.month == 1:
            tasks.append("yearly")

        if tasks:
            async with session_factory() as db:
                result = await db.execute(select(Project).where(Project.status == ProjectStatus.ACTIVE))
                projects = result.scalars().all()

            for report_type in tasks:
                for project in projects:
                    try:
                        await generate_report(session_factory, project.id, report_type)
                    except Exception as e:
                        logger.error(f"Scheduled {report_type} report failed for project {project.id}: {e}")

        await asyncio.sleep(300)  # check every 5 minutes


async def main():
    if not TEMPLATE_DIR.exists():
        logger.warning("Template directory not found: %s", TEMPLATE_DIR)

    shutdown_event = asyncio.Event()
    loop = asyncio.get_running_loop()
    for sig in (signal.SIGTERM, signal.SIGINT):
        loop.add_signal_handler(sig, shutdown_event.set)

    engine, session_factory = create_db(DATABASE_URL)
    bus = EventBus(REDIS_URL)
    await bus.connect()

    liveness_port = int(os.getenv("LIVENESS_PORT", "9097"))
    liveness = await start_liveness_server(liveness_port)

    logger.info("Report Service started")
    try:
        await asyncio.gather(
            listen_for_requests(bus, session_factory, shutdown_event),
            scheduled_reports(bus, session_factory, shutdown_event),
        )
    finally:
        liveness.close()
        logger.info("Report Service shutting down, metrics: %s", metrics.summary())
        await bus.close()
        await engine.dispose()


if __name__ == "__main__":
    asyncio.run(main())
